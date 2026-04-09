from __future__ import annotations

from pathlib import Path
from types import SimpleNamespace

import fitz
import pytest
from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation

from mcp_servers.file_tools.utils import file_extractors
from mcp_servers.file_tools.utils.file_extractors import extract_preview


def _fake_settings(*, is_multimodal: bool) -> SimpleNamespace:
    return SimpleNamespace(
        llm=SimpleNamespace(is_multimodal=is_multimodal),
        preview=SimpleNamespace(
            text_depth_chars={"L1": 2000, "L2": 5000, "L3": 8000},
            excel_depth_rows={"L1": 10, "L2": 50, "L3": 100},
            pdf_depth_pages={"L1": 1, "L2": 2, "L3": 3},
            pdf_render_scale=3.0,
            image_max_edge=1600,
        ),
    )


def test_extract_preview_for_text_file(tmp_path: Path) -> None:
    file_path = tmp_path / "notes.md"
    file_path.write_text("第一行\n第二行\n第三行", encoding="utf-8")

    preview = extract_preview(file_path, "L1")

    assert preview.file_type == "md"
    assert "第二行" in preview.preview_text
    assert preview.metadata["line_count"] == 3


def test_extract_preview_for_docx_file(tmp_path: Path) -> None:
    file_path = tmp_path / "report.docx"
    document = Document()
    document.add_paragraph("项目周报")
    document.add_paragraph("本周完成搜索工具联调。")
    document.save(file_path)

    preview = extract_preview(file_path, "L2")

    assert preview.file_type == "docx"
    assert "项目周报" in preview.preview_text
    assert preview.metadata["paragraph_count"] == 2


def test_extract_preview_for_xlsx_file(tmp_path: Path) -> None:
    file_path = tmp_path / "data.xlsx"
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Sheet1"
    worksheet.append(["name", "score"])
    worksheet.append(["Flora", 95])
    workbook.save(file_path)

    preview = extract_preview(file_path, "L1")

    assert preview.file_type == "xlsx"
    assert "[Sheet] Sheet1" in preview.preview_text
    assert preview.metadata["sheet_count"] == 1
    assert preview.metadata["sheet_names"] == ["Sheet1"]


def test_extract_preview_for_pptx_file(tmp_path: Path) -> None:
    file_path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Transformer 讲解"
    slide.placeholders[1].text = "第 2 页有完整架构图"
    presentation.save(file_path)

    preview = extract_preview(file_path, "L1")

    assert preview.file_type == "pptx"
    assert "Transformer 讲解" in preview.preview_text
    assert preview.metadata["slide_count"] == 1


def test_extract_preview_for_pdf_file_prefers_visual_mode(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    file_path = tmp_path / "paper.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "LoRA 可以显著减少训练显存占用")
    document.save(file_path)
    document.close()

    def fake_visual_summary(assets, *, depth: str, file_name: str, file_type: str) -> str:
        assert len(assets) == 1
        assert depth == "L1"
        assert file_type == "pdf"
        return "这是一个关于 LoRA 的 PDF 页面预览摘要。"

    monkeypatch.setattr(file_extractors, "get_settings", lambda: _fake_settings(is_multimodal=True))
    monkeypatch.setattr(file_extractors, "summarize_visual_assets", fake_visual_summary)

    preview = extract_preview(file_path, "L1")

    assert preview.file_type == "pdf"
    assert "LoRA" in preview.preview_text
    assert preview.metadata["preview_mode"] == "pdf_rendered_image_summary"
    assert preview.metadata["preview_pages"] == 1


def test_extract_preview_for_pdf_file_falls_back_to_text_when_visual_fails(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    file_path = tmp_path / "paper.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "LoRA 可以显著减少训练显存占用")
    document.save(file_path)
    document.close()

    def raise_visual_error(*args, **kwargs) -> str:
        raise RuntimeError("vision unavailable")

    monkeypatch.setattr(file_extractors, "get_settings", lambda: _fake_settings(is_multimodal=True))
    monkeypatch.setattr(file_extractors, "summarize_visual_assets", raise_visual_error)

    preview = extract_preview(file_path, "L1")

    assert preview.file_type == "pdf"
    assert "视觉模型查看失败" in preview.preview_text
    assert "LoRA" in preview.preview_text
    assert preview.metadata["preview_mode"] == "pdf_text_fallback"
    assert "vision unavailable" in str(preview.metadata["fallback_reason"])


def test_extract_preview_for_pdf_file_uses_text_preview_when_multimodal_disabled(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    file_path = tmp_path / "paper.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "latest project report")
    document.save(file_path)
    document.close()

    def fail_if_visual_called(*args, **kwargs) -> str:
        raise AssertionError("visual summary should not be called")

    monkeypatch.setattr(file_extractors, "get_settings", lambda: _fake_settings(is_multimodal=False))
    monkeypatch.setattr(file_extractors, "summarize_visual_assets", fail_if_visual_called)

    preview = extract_preview(file_path, "L1")

    assert preview.file_type == "pdf"
    assert "当前模型未启用多模态预览" in preview.preview_text
    assert "latest project report" in preview.preview_text
    assert preview.metadata["preview_mode"] == "pdf_text_preview"
    assert preview.metadata["fallback_reason"] == "vision_disabled_by_config"


def test_extract_preview_for_image_file_uses_visual_summary(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    file_path = tmp_path / "chart.png"
    image = Image.new("RGB", (640, 480), color="white")
    image.save(file_path)

    def fake_visual_summary(assets, *, depth: str, file_name: str, file_type: str) -> str:
        assert len(assets) == 1
        assert file_type == "png"
        return "这是一张白底图表图片。"

    monkeypatch.setattr(file_extractors, "get_settings", lambda: _fake_settings(is_multimodal=True))
    monkeypatch.setattr(file_extractors, "summarize_visual_assets", fake_visual_summary)

    preview = extract_preview(file_path, "L1")

    assert preview.file_type == "png"
    assert "白底图表" in preview.preview_text
    assert preview.metadata["width"] == 640
    assert preview.metadata["height"] == 480
    assert preview.metadata["preview_mode"] == "image_visual_summary"


def test_extract_preview_for_image_file_falls_back_to_metadata(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    file_path = tmp_path / "chart.png"
    image = Image.new("RGB", (640, 480), color="white")
    image.save(file_path)

    def raise_visual_error(*args, **kwargs) -> str:
        raise RuntimeError("vision unavailable")

    monkeypatch.setattr(file_extractors, "get_settings", lambda: _fake_settings(is_multimodal=True))
    monkeypatch.setattr(file_extractors, "summarize_visual_assets", raise_visual_error)

    preview = extract_preview(file_path, "L1")

    assert preview.file_type == "png"
    assert "视觉模型查看失败" in preview.preview_text
    assert preview.metadata["preview_mode"] == "image_metadata_fallback"


def test_extract_preview_for_image_file_uses_metadata_when_multimodal_disabled(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    file_path = tmp_path / "chart.png"
    image = Image.new("RGB", (640, 480), color="white")
    image.save(file_path)

    def fail_if_visual_called(*args, **kwargs) -> str:
        raise AssertionError("visual summary should not be called")

    monkeypatch.setattr(file_extractors, "get_settings", lambda: _fake_settings(is_multimodal=False))
    monkeypatch.setattr(file_extractors, "summarize_visual_assets", fail_if_visual_called)

    preview = extract_preview(file_path, "L1")

    assert preview.file_type == "png"
    assert "当前模型未启用多模态预览" in preview.preview_text
    assert preview.metadata["preview_mode"] == "image_metadata_fallback"
    assert preview.metadata["fallback_reason"] == "vision_disabled_by_config"


def test_extract_preview_for_zero_kb_pdf_raises_structured_error(tmp_path: Path) -> None:
    file_path = tmp_path / "empty.pdf"
    file_path.write_bytes(b"")

    with pytest.raises(file_extractors.StructuredPreviewError) as exc_info:
        extract_preview(file_path, "L1")

    assert exc_info.value.error_type == "pdf_empty_file"


def test_extract_preview_for_blank_pdf_raises_structured_error(tmp_path: Path) -> None:
    file_path = tmp_path / "blank.pdf"
    document = fitz.open()
    document.new_page()
    document.save(file_path)
    document.close()

    monkeypatch = pytest.MonkeyPatch()
    monkeypatch.setattr(file_extractors, "get_settings", lambda: _fake_settings(is_multimodal=True))

    with pytest.raises(file_extractors.StructuredPreviewError) as exc_info:
        extract_preview(file_path, "L1")

    assert exc_info.value.error_type == "pdf_no_visual_content"
    monkeypatch.undo()
