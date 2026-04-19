from __future__ import annotations

from dataclasses import dataclass
from decimal import ROUND_HALF_UP, Decimal
from io import BytesIO
from pathlib import Path

import fitz
from docx import Document
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation

from config.settings import get_settings
from mcp_servers.file_tools.utils.time_format import format_modified_timestamp

TEXT_FILE_SUFFIXES = {
    ".csv",
    ".json",
    ".log",
    ".md",
    ".py",
    ".txt",
    ".yaml",
    ".yml",
}
IMAGE_FILE_SUFFIXES = {".bmp", ".gif", ".jpeg", ".jpg", ".png", ".webp"}
SUPPORTED_DEPTHS = {"L1", "L2", "L3"}
BLANK_IMAGE_TOLERANCE = 2


class UnsupportedFileTypeError(ValueError):
    pass


class StructuredPreviewError(RuntimeError):
    def __init__(
        self,
        error_type: str,
        message: str,
        *,
        retryable: bool = False,
        user_hint: str | None = None,
        operator_hint: str | None = None,
    ) -> None:
        super().__init__(message)
        self.error_type = error_type
        self.message = message
        self.retryable = retryable
        self.user_hint = user_hint or message
        self.operator_hint = operator_hint or ""


@dataclass(frozen=True)
class FilePreview:
    file_type: str
    preview_text: str
    metadata: dict[str, object]


@dataclass(frozen=True)
class PdfVisualSource:
    metadata: dict[str, object]


def normalize_depth(depth: str) -> str:
    normalized = depth.strip().upper()
    if normalized not in SUPPORTED_DEPTHS:
        raise ValueError("depth 仅支持 L1、L2、L3。")
    return normalized


def _format_decimal_size(size: int, unit_base: int) -> str:
    value = Decimal(size) / Decimal(unit_base)
    rounded = value.quantize(Decimal("0.1"), rounding=ROUND_HALF_UP)
    return f"{rounded}"


def format_size_display(size: int) -> str:
    if size < 1024:
        return f"{size} B"
    if size < 1024**2:
        return f"{_format_decimal_size(size, 1024)} KB"
    if size < 1024**3:
        return f"{_format_decimal_size(size, 1024**2)} MB"
    return f"{_format_decimal_size(size, 1024**3)} GB"


def extract_preview(file_path: Path, depth: str) -> FilePreview:
    normalized_depth = normalize_depth(depth)
    suffix = file_path.suffix.lower()

    if suffix in TEXT_FILE_SUFFIXES:
        return _extract_text_preview(file_path, normalized_depth, file_type=suffix.lstrip("."))
    if suffix == ".docx":
        return _extract_docx_preview(file_path, normalized_depth)
    if suffix == ".xlsx":
        return _extract_xlsx_preview(file_path, normalized_depth)
    if suffix == ".pptx":
        return _extract_pptx_preview(file_path, normalized_depth)
    if suffix == ".pdf":
        return _extract_pdf_preview(file_path, normalized_depth)
    if suffix in IMAGE_FILE_SUFFIXES:
        return _extract_image_preview(file_path, normalized_depth)

    raise UnsupportedFileTypeError(f"暂不支持读取 {suffix or '无扩展名'} 文件。")


def build_common_metadata(file_path: Path) -> dict[str, object]:
    stat_result = file_path.stat()
    return {
        "size": stat_result.st_size,
        "size_display": format_size_display(stat_result.st_size),
        "modified": format_modified_timestamp(stat_result.st_mtime),
    }


def _text_char_limit(depth: str) -> int:
    return get_settings().preview.text_depth_chars[depth]


def _excel_row_limit(depth: str) -> int:
    return get_settings().preview.excel_depth_rows[depth]


def _pdf_page_limit(depth: str) -> int:
    return get_settings().preview.pdf_depth_pages[depth]


def _is_multimodal_enabled() -> bool:
    return bool(get_settings().llm.is_multimodal)


def _truncate_text(text: str, limit: int) -> tuple[str, bool]:
    normalized_text = text.replace("\r\n", "\n").replace("\r", "\n").strip()
    if len(normalized_text) <= limit:
        return normalized_text, False
    return normalized_text[:limit].rstrip() + "...", True


def _merge_metadata(file_path: Path, extra_metadata: dict[str, object]) -> dict[str, object]:
    metadata = build_common_metadata(file_path)
    metadata.update(extra_metadata)
    return metadata


def _read_text_content(file_path: Path) -> str:
    try:
        return file_path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return file_path.read_text(encoding="utf-8", errors="replace")


def _extract_text_preview(file_path: Path, depth: str, *, file_type: str) -> FilePreview:
    content = _read_text_content(file_path)
    preview_text, truncated = _truncate_text(content, _text_char_limit(depth))
    metadata = _merge_metadata(
        file_path,
        {
            "line_count": len(content.splitlines()),
            "char_count": len(content),
            "truncated": truncated,
        },
    )
    return FilePreview(file_type=file_type, preview_text=preview_text, metadata=metadata)


def _extract_docx_preview(file_path: Path, depth: str) -> FilePreview:
    document = Document(file_path)
    paragraphs = [
        paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()
    ]
    content = "\n".join(paragraphs)
    preview_text, truncated = _truncate_text(content, _text_char_limit(depth))
    metadata = _merge_metadata(
        file_path,
        {
            "paragraph_count": len(paragraphs),
            "char_count": len(content),
            "truncated": truncated,
        },
    )
    return FilePreview(file_type="docx", preview_text=preview_text, metadata=metadata)


def _extract_xlsx_preview(file_path: Path, depth: str) -> FilePreview:
    workbook = load_workbook(file_path, read_only=True, data_only=True)
    row_limit = _excel_row_limit(depth)
    sheet_names: list[str] = []
    preview_sections: list[str] = []

    for worksheet in workbook.worksheets:
        sheet_names.append(worksheet.title)
        preview_sections.append(f"[Sheet] {worksheet.title}")
        for row_index, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
            if row_index > row_limit:
                break
            values = [str(cell) if cell is not None else "" for cell in row]
            preview_sections.append(" | ".join(values).rstrip())

    content = "\n".join(section for section in preview_sections if section)
    preview_text, truncated = _truncate_text(content, _text_char_limit(depth))
    metadata = _merge_metadata(
        file_path,
        {
            "sheet_count": len(sheet_names),
            "sheet_names": sheet_names,
            "preview_rows_per_sheet": row_limit,
            "truncated": truncated,
        },
    )
    workbook.close()
    return FilePreview(file_type="xlsx", preview_text=preview_text, metadata=metadata)


def _extract_pptx_preview(file_path: Path, depth: str) -> FilePreview:
    presentation = Presentation(file_path)
    slide_limit = _pdf_page_limit(depth)
    preview_sections: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        if slide_index > slide_limit:
            break
        shape_texts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                shape_texts.append(text.strip())
        preview_sections.append(f"[Slide {slide_index}]")
        preview_sections.append("\n".join(shape_texts) if shape_texts else "<无可提取文本>")

    content = "\n".join(preview_sections)
    preview_text, truncated = _truncate_text(content, _text_char_limit(depth))
    metadata = _merge_metadata(
        file_path,
        {
            "slide_count": len(presentation.slides),
            "preview_slides": min(len(presentation.slides), slide_limit),
            "truncated": truncated,
        },
    )
    return FilePreview(file_type="pptx", preview_text=preview_text, metadata=metadata)


def _extract_pdf_preview(file_path: Path, depth: str) -> FilePreview:
    is_multimodal = _is_multimodal_enabled()
    try:
        pdf_source = _collect_pdf_visual_source(file_path, depth)
    except StructuredPreviewError:
        raise
    except Exception as exc:
        return _build_pdf_text_preview_result(
            file_path,
            depth,
            preview_mode="pdf_text_source",
            fallback_reason="vision_disabled_by_config" if not is_multimodal else str(exc),
        )

    if not is_multimodal:
        return _build_pdf_text_preview_result(
            file_path,
            depth,
            preview_mode="pdf_text_source",
            fallback_reason="vision_disabled_by_config",
            base_metadata=pdf_source.metadata,
        )

    preview = _extract_pdf_text_preview(file_path, depth)
    metadata = dict(preview.metadata)
    metadata.update(pdf_source.metadata)
    metadata["preview_mode"] = "pdf_visual_source"
    return FilePreview(file_type="pdf", preview_text=preview.preview_text, metadata=metadata)


def _extract_pdf_text_preview(file_path: Path, depth: str) -> FilePreview:
    document = _open_pdf_document(file_path)
    page_limit = _pdf_page_limit(depth)
    preview_sections: list[str] = []

    for page_index in range(min(document.page_count, page_limit)):
        page = document.load_page(page_index)
        preview_sections.append(f"[Page {page_index + 1}]")
        preview_sections.append(page.get_text("text").strip() or "<当前页可提取文本较少>")

    content = "\n".join(preview_sections)
    preview_text, truncated = _truncate_text(content, _text_char_limit(depth))
    metadata = _merge_metadata(
        file_path,
        {
            "page_count": document.page_count,
            "preview_pages": min(document.page_count, page_limit),
            "truncated": truncated,
        },
    )
    document.close()
    return FilePreview(file_type="pdf", preview_text=preview_text, metadata=metadata)


def _extract_image_preview(file_path: Path, depth: str) -> FilePreview:
    del depth
    file_type = file_path.suffix.lower().lstrip(".")
    base_metadata = _build_image_metadata(file_path)
    is_multimodal = _is_multimodal_enabled()

    if is_multimodal:
        return _build_image_metadata_preview(
            file_path,
            file_type=file_type,
            base_metadata=base_metadata,
            preview_mode="image_visual_source",
            fallback_reason=None,
            notice=None,
        )

    return _build_image_metadata_preview(
        file_path,
        file_type=file_type,
        base_metadata=base_metadata,
        preview_mode="image_metadata_source",
        fallback_reason="vision_disabled_by_config",
        notice="当前模型未启用多模态预览。",
    )


def _build_image_metadata(file_path: Path) -> dict[str, object]:
    with Image.open(file_path) as image:
        image.load()
        return {
            "format": image.format or file_path.suffix.lower().lstrip("."),
            "width": image.width,
            "height": image.height,
            "mode": image.mode,
        }


def _collect_pdf_visual_source(file_path: Path, depth: str) -> PdfVisualSource:
    document = _open_pdf_document(file_path)
    try:
        page_limit = min(document.page_count, _pdf_page_limit(depth))
        render_scale = get_settings().preview.pdf_render_scale
        rendered_pages = 0
        blank_pages: list[int] = []

        for page_index in range(page_limit):
            page = document.load_page(page_index)
            pixmap = page.get_pixmap(matrix=fitz.Matrix(render_scale, render_scale), alpha=False)
            if pixmap.width <= 0 or pixmap.height <= 0:
                raise StructuredPreviewError(
                    "pdf_blank_render",
                    "当前 PDF 页面渲染结果为空，无法生成图片预览。",
                    user_hint="这个 PDF 的页面无法正常渲染，请尝试更换文件。",
                    operator_hint="PDF 页渲染得到的图像宽高无效。",
                )

            image_bytes = pixmap.tobytes("png")
            if not image_bytes:
                raise StructuredPreviewError(
                    "pdf_blank_render",
                    "当前 PDF 页面渲染结果为空，无法生成图片预览。",
                    user_hint="这个 PDF 的页面无法正常渲染，请尝试更换文件。",
                    operator_hint="pixmap.tobytes('png') 返回空字节。",
                )

            rendered_pages += 1
            if _is_effectively_blank_image(image_bytes):
                blank_pages.append(page_index + 1)

        if rendered_pages == 0:
            raise StructuredPreviewError(
                "pdf_blank_render",
                "当前 PDF 页面渲染结果为空，无法生成图片预览。",
                user_hint="这个 PDF 的页面无法正常渲染，请尝试更换文件。",
                operator_hint="PDF 渲染流程结束后没有生成任何视觉资产。",
            )

        if len(blank_pages) == rendered_pages:
            raise StructuredPreviewError(
                "pdf_no_visual_content",
                "当前 PDF 渲染后的页面没有检测到有效视觉内容。",
                user_hint="这个 PDF 可能是空白页、纯黑/纯白页，暂时无法预览。",
                operator_hint=f"预览页全部被判定为近似纯黑/纯白：pages={blank_pages}。",
            )

        return PdfVisualSource(
            metadata={
                "page_count": document.page_count,
                "preview_pages": rendered_pages,
                "render_scale": render_scale,
            },
        )
    finally:
        document.close()


def _build_pdf_text_preview_result(
    file_path: Path,
    depth: str,
    *,
    preview_mode: str,
    fallback_reason: str | None = None,
    base_metadata: dict[str, object] | None = None,
) -> FilePreview:
    preview = _extract_pdf_text_preview(file_path, depth)
    metadata = dict(base_metadata or {})
    metadata.update(preview.metadata)
    metadata["preview_mode"] = preview_mode
    if fallback_reason is not None:
        metadata["fallback_reason"] = fallback_reason
    return FilePreview(
        file_type="pdf",
        preview_text=preview.preview_text,
        metadata=metadata,
    )


def _build_image_metadata_preview(
    file_path: Path,
    *,
    file_type: str,
    base_metadata: dict[str, object],
    preview_mode: str,
    fallback_reason: str | None,
    notice: str | None,
) -> FilePreview:
    metadata = dict(base_metadata)
    metadata["preview_mode"] = preview_mode
    if fallback_reason is not None:
        metadata["fallback_reason"] = fallback_reason
    return FilePreview(
        file_type=file_type,
        preview_text=_build_image_metadata_fallback_text(base_metadata, notice=notice),
        metadata=_merge_metadata(file_path, metadata),
    )


def _build_image_metadata_fallback_text(
    metadata: dict[str, object], *, notice: str | None = None
) -> str:
    image_format = metadata.get("format", "image")
    width = metadata.get("width", "?")
    height = metadata.get("height", "?")
    if notice:
        return f"{notice} 当前返回基础信息：格式 {image_format}，尺寸 {width}x{height}。"
    return f"图片基础信息：格式 {image_format}，尺寸 {width}x{height}。"


def _open_pdf_document(file_path: Path) -> fitz.Document:
    file_size = file_path.stat().st_size
    if file_size == 0:
        raise StructuredPreviewError(
            "pdf_empty_file",
            "当前 PDF 文件大小为 0KB，无法生成预览。",
            user_hint="这个 PDF 是空文件，请更换文件或重新导出。",
            operator_hint="检测到 PDF 文件大小为 0 字节，直接拦截预览。",
        )

    document = fitz.open(file_path)
    if document.page_count <= 0:
        document.close()
        raise StructuredPreviewError(
            "pdf_no_pages",
            "当前 PDF 不包含可预览的页面。",
            user_hint="这个 PDF 没有可预览页面，请确认文件是否损坏。",
            operator_hint="fitz 打开 PDF 成功，但 page_count <= 0。",
        )
    return document


def _is_effectively_blank_image(image_bytes: bytes) -> bool:
    with Image.open(BytesIO(image_bytes)) as image:
        grayscale = image.convert("L")
        extrema = grayscale.getextrema()
    if extrema is None:
        return True
    low, high = extrema
    if low == high and low in {0, 255}:
        return True
    return (high - low) <= BLANK_IMAGE_TOLERANCE and (
        low <= BLANK_IMAGE_TOLERANCE or high >= 255 - BLANK_IMAGE_TOLERANCE
    )
